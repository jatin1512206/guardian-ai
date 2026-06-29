import sys
import subprocess

# Auto-install python-pptx if missing
try:
    import pptx
except ImportError:
    print("Installing python-pptx for PowerPoint generation...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation_pptx(output_filename="guardian_ai_presentation.pptx"):
    prs = Presentation()
    
    # Set to widescreen 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Clean high-contrast palette
    WHITE_BG = RGBColor(255, 255, 255)      # Clean white slide backgrounds
    NAVY = RGBColor(15, 23, 42)             # Deep Navy (#0f172a) for main headers
    TEAL = RGBColor(14, 116, 144)           # #0e7490 for subtitles/accents
    CHARCOAL = RGBColor(51, 65, 85)         # #334155 for body text
    LIGHT_GRAY = RGBColor(241, 245, 249)    # #f1f5f9 for card background panels
    BORDER_GRAY = RGBColor(226, 232, 240)   # #e2e8f0 for borders

    blank_layout = prs.slide_layouts[6]     # Blank layout to manually position elements

    def apply_slide_background(slide):
        # Set slide background to clean solid white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE_BG

    def add_slide_header(slide, title_text):
        apply_slide_background(slide)
        
        # Draw a beautiful dark navy header block at the top of the slide
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.333)
        height = Inches(1.1)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background() # No border
        
        # Add Title text directly as a run inside the header block text frame
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12.133), Inches(0.7))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        run = p.add_run()
        run.text = title_text
        run.font.name = "Arial"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE_BG

    # ================= Slide 1: Title Slide =================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)

    # Accent color block on the left (Teal)
    accent_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAL
    accent_bar.line.fill.background()

    # Main Title Box
    tx_title = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(1.8))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    
    # Title Text Run
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = "GUARDIAN-AI 🛡️"
    run_title.font.name = "Arial"
    run_title.font.size = Pt(56)
    run_title.font.bold = True
    run_title.font.color.rgb = NAVY

    # Subtitle Text Run
    p_sub = tf_title.add_paragraph()
    p_sub.alignment = PP_ALIGN.LEFT
    p_sub.space_before = Pt(12)
    run_sub = p_sub.add_run()
    run_sub.text = "Edge-Computer Vision & Asynchronous Multi-Agent Safety Dashboard"
    run_sub.font.name = "Arial"
    run_sub.font.size = Pt(22)
    run_sub.font.bold = True
    run_sub.font.color.rgb = TEAL

    # Meta/Author info
    tx_meta = slide1.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.5), Inches(1.2))
    tf_meta = tx_meta.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    
    run_meta1 = p_meta.add_run()
    run_meta1.text = "Submission Proposal: Hack2Skill Data & AI Challenge\n"
    run_meta1.font.name = "Arial"
    run_meta1.font.size = Pt(14)
    run_meta1.font.bold = True
    run_meta1.font.color.rgb = CHARCOAL

    run_meta2 = p_meta.add_run()
    run_meta2.text = "Developer & Architect: Jatin Rajpoot"
    run_meta2.font.name = "Arial"
    run_meta2.font.size = Pt(13)
    run_meta2.font.color.rgb = TEAL

    # ================= Slide 2: The Problem =================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "The Problem: Cognitive Distractions & Road Hazards")

    tx_content = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf2 = tx_content.text_frame
    tf2.word_wrap = True

    # Section description
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = "Critical Driver Safety Gaps:"
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = TEAL
    p.space_after = Pt(10)

    p_desc = tf2.add_paragraph()
    run_desc = p_desc.add_run()
    run_desc.text = "Over 1.3 million global vehicular deaths occur annually due to human error, driver fatigue, and phone distractions. Current dashboard recorders act as passive blackboxes, lacking the intelligence to predict accidents and alert drivers before impact."
    run_desc.font.name = "Arial"
    run_desc.font.size = Pt(15)
    run_desc.font.color.rgb = CHARCOAL
    p_desc.space_after = Pt(20)

    # Bullet points with styled runs
    bullets = [
        ("Driver Fatigue & Micro-sleeps", "Split-second closures of eyes lead directly to road departures and high-speed crashes before drivers realize they are tired."),
        ("Visual Distractions (Smartphone Usage)", "Taking hands off the wheel or shifting eye gaze to screens dramatically increases stopping distance."),
        ("Passive Dashcam Monitoring", "Traditional cameras lack edge computer vision to compute real-time risks and alert drivers actively.")
    ]

    for title, desc in bullets:
        p_b = tf2.add_paragraph()
        p_b.space_before = Pt(8)
        
        run_title = p_b.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.name = "Arial"
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = TEAL

        run_text = p_b.add_run()
        run_text.text = desc
        run_text.font.name = "Arial"
        run_text.font.size = Pt(14)
        run_text.font.color.rgb = CHARCOAL

    # ================= Slide 3: What We Built =================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "What We Built: Closed-Loop AI Safety System")

    # 4 distinct visual cards
    cards = [
        ("🛡️ Browser AI CV", "MediaPipe Holistic tracks 3D face mesh targets, arm segments, and 21-joint finger skeletons in the browser.", TEAL),
        ("🧠 Asynchronous Agents", "FastAPI agents (Driver, Vehicle, Predictor, Intervention) process telemetry events on an async bus.", NAVY),
        ("🔊 Synth Siren Alerts", "Web Audio oscillators produce urgent sirens and warning chimes during safety violations.", TEAL),
        ("🎛️ Interactive Controls", "Speed, steering, and lane sliders allow manual overrides to instantly test accident risks.", NAVY)
    ]

    for i, (title, desc, color) in enumerate(cards):
        left = Inches(0.6 + i * 3.05)
        top = Inches(1.8)
        width = Inches(2.9)
        height = Inches(4.7)

        # Draw card panel background
        card_panel = slide3.shapes.add_shape(1, left, top, width, height)
        card_panel.fill.solid()
        card_panel.fill.fore_color.rgb = LIGHT_GRAY
        card_panel.line.color.rgb = BORDER_GRAY
        card_panel.line.width = Pt(1.5)

        # Add card text box
        tx = slide3.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.4))
        tf_card = tx.text_frame
        tf_card.word_wrap = True

        p_t = tf_card.paragraphs[0]
        run_t = p_t.add_run()
        run_t.text = title
        run_t.font.name = "Arial"
        run_t.font.size = Pt(18)
        run_t.font.bold = True
        run_t.font.color.rgb = color
        p_t.space_after = Pt(15)

        p_d = tf_card.add_paragraph()
        run_d = p_d.add_run()
        run_d.text = desc
        run_d.font.name = "Arial"
        run_d.font.size = Pt(13)
        run_d.font.color.rgb = CHARCOAL
        p_d.space_before = Pt(10)

    # ================= Slide 4: Why We Built It This Way =================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "Why We Built It That Way: Design Rationale")

    tx_why = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf_why = tx_why.text_frame
    tf_why.word_wrap = True

    why_points = [
        ("Edge-CV + Cloud Fusion Model", "MediaPipe Holistic runs client-side in the browser. This eliminates expensive cloud GPU server hosting costs and guarantees zero latency overhead for frame capture processing (<30ms)."),
        ("Asynchronous Event Bus Decoupling", "Decouples telemetry loops (vehicle updates at 10Hz, camera landmarks at 30Hz). If the camera is toggled off, backend agents fallback to physics-simulations instantly."),
        ("Direct DOM Telemetry Updates", "Rendering 30 FPS face mesh coords inside React state freezes the virtual DOM. We update HUD metrics directly in the DOM by ID, saving significant CPU cycles.")
    ]

    for title, desc in why_points:
        p_title = tf_why.add_paragraph()
        p_title.space_before = Pt(12)
        run_t = p_title.add_run()
        run_t.text = f"✔  {title}: "
        run_t.font.name = "Arial"
        run_t.font.size = Pt(16)
        run_t.font.bold = True
        run_t.font.color.rgb = TEAL

        p_desc = tf_why.add_paragraph()
        p_desc.space_before = Pt(2)
        p_desc.space_after = Pt(8)
        run_d = p_desc.add_run()
        run_d.text = desc
        run_d.font.name = "Arial"
        run_d.font.size = Pt(14)
        run_d.font.color.rgb = CHARCOAL

    # ================= Slide 5: How It Works =================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "How It Works: Technical Data Flow")

    tx_flow = slide5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf_flow = tx_flow.text_frame
    tf_flow.word_wrap = True

    flow_steps = [
        ("1. Sensation", "Webcam captures raw frame video. MediaPipe tracks Eye Aspect Ratio (EAR), Mouth aspect (MAR), and Hand-to-Cheek distance vectors."),
        ("2. Transmission", "WebSockets stream coordinates to FastAPI's async event bus, throttled to 150ms to prevent server buffer overflows."),
        ("3. Sequence Modeling", "Vehicle telemetry sequence (50, 10) passes through a BiGRU network; driver parameters are fused through a Transformer cross-attention encoder."),
        ("4. Forecast", "LSTM forecaster outputs collision probability and type (e.g. rear-end) 3-5 seconds ahead."),
        ("5. Actuation", "Intervention Agent triggers wailing sirens in the browser and logs safety event audits into SQLite WAL database.")
    ]

    for step_title, step_desc in flow_steps:
        p_step = tf_flow.add_paragraph()
        p_step.space_before = Pt(10)
        
        run_st = p_step.add_run()
        run_st.text = f"•  {step_title}: "
        run_st.font.name = "Arial"
        run_st.font.size = Pt(15)
        run_st.font.bold = True
        run_st.font.color.rgb = TEAL

        run_sd = p_step.add_run()
        run_sd.text = step_desc
        run_sd.font.name = "Arial"
        run_sd.font.size = Pt(14)
        run_sd.font.color.rgb = CHARCOAL

    # ================= Slide 6: Model Training & Verification =================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "Trained Models & System Verification")

    tx_verif = slide6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    tf_verif = tx_verif.text_frame
    tf_verif.word_wrap = True

    p = tf_verif.paragraphs[0]
    run_h = p.add_run()
    run_h.text = "1. Active PyTorch Neural Network Weights:"
    run_h.font.name = "Arial"
    run_h.font.size = Pt(16)
    run_h.font.bold = True
    run_h.font.color.rgb = NAVY
    p.space_after = Pt(8)

    models_list = [
        ("Driver Monitor", "MobileNetV2 CNN Classifier", "driver_monitor.pt"),
        ("Vehicle Dynamics", "Bi-GRU with Attention", "vehicle_behavior.pt"),
        ("Sensor Fusion", "Transformer Cross-Attention", "sensor_fusion.pt"),
        ("Accident Predictor", "LSTM Sequence Forecaster", "accident_predictor.pt")
    ]

    for model, arch, filename in models_list:
        p_m = tf_verif.add_paragraph()
        p_m.space_after = Pt(2)
        
        run_lbl = p_m.add_run()
        run_lbl.text = f"  - {model} ({arch})  -->  "
        run_lbl.font.name = "Arial"
        run_lbl.font.size = Pt(14)
        run_lbl.font.color.rgb = CHARCOAL
        
        run_file = p_m.add_run()
        run_file.text = filename
        run_file.font.name = "Arial"
        run_file.font.size = Pt(14)
        run_file.font.bold = True
        run_file.font.color.rgb = TEAL

    p_t = tf_verif.add_paragraph()
    p_t.space_before = Pt(18)
    run_th = p_t.add_run()
    run_th.text = "2. System Verification Metrics:"
    run_th.font.name = "Arial"
    run_th.font.size = Pt(16)
    run_th.font.bold = True
    run_th.font.color.rgb = NAVY
    p_t.space_after = Pt(6)

    p_td = tf_verif.add_paragraph()
    run_td1 = p_td.add_run()
    run_td1.text = "✔  Pytest Suite: "
    run_td1.font.name = "Arial"
    run_td1.font.size = Pt(14)
    run_td1.font.bold = True
    run_td1.font.color.rgb = TEAL
    
    run_td2 = p_td.add_run()
    run_td2.text = "10/10 tests pass, validating async queues, DB persistence, API, and sequence networks.\n"
    run_td2.font.name = "Arial"
    run_td2.font.size = Pt(14)
    run_td2.font.color.rgb = CHARCOAL

    run_td3 = p_td.add_run()
    run_td3.text = "✔  Inference Latency: "
    run_td3.font.name = "Arial"
    run_td3.font.size = Pt(14)
    run_td3.font.bold = True
    run_td3.font.color.rgb = TEAL

    run_td4 = p_td.add_run()
    run_td4.text = "Inference is completed under 10ms via PyTorch cpu optimization layers, exceeding the 100ms budget."
    run_td4.font.name = "Arial"
    run_td4.font.size = Pt(14)
    run_td4.font.color.rgb = CHARCOAL

    # Save presentation
    prs.save(output_filename)
    print("PowerPoint Presentation successfully generated: guardian_ai_presentation.pptx")

if __name__ == "__main__":
    create_presentation_pptx()
